"""Validate generated document containers without relying on pixel snapshots.

The checks are intentionally deterministic: container validity, page/slide
geometry, readable text presence and absence of internal placeholders. Dynamic
AI prose is not compared byte-for-byte.
"""

from __future__ import annotations

import argparse
import json
import re
import zipfile
from pathlib import Path
from typing import Any


PLACEHOLDER_PATTERNS = (
    re.compile(r"\bTODO\b", re.IGNORECASE),
    re.compile(r"\{\{[^}]+\}\}"),
    re.compile(r"```(?:json|markdown|python)?", re.IGNORECASE),
    re.compile(r"internal ai|system prompt|prompt injection", re.IGNORECASE),
)


def _placeholder_hits(text: str) -> list[str]:
    return [pattern.pattern for pattern in PLACEHOLDER_PATTERNS if pattern.search(text)]


def validate_pdf(path: Path) -> dict[str, Any]:
    from pypdf import PdfReader

    reader = PdfReader(str(path), strict=False)
    errors: list[str] = []
    if reader.is_encrypted:
        errors.append("PDF-encrypted")
    if not reader.pages:
        errors.append("PDF-without-pages")
    text_chars = 0
    for index, page in enumerate(reader.pages, start=1):
        box = page.mediabox
        if float(box.width) <= 0 or float(box.height) <= 0:
            errors.append(f"page-{index}-invalid-geometry")
        text = page.extract_text() or ""
        text_chars += len(text.strip())
        errors.extend(f"page-{index}-placeholder:{hit}" for hit in _placeholder_hits(text))
    if text_chars == 0:
        errors.append("PDF-without-extractable-text")
    return {"format": "pdf", "pages": len(reader.pages), "text_chars": text_chars, "errors": errors}


def validate_pptx(path: Path) -> dict[str, Any]:
    from pptx import Presentation

    presentation = Presentation(str(path))
    errors: list[str] = []
    slide_width = presentation.slide_width
    slide_height = presentation.slide_height
    text_chars = 0
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape_number, shape in enumerate(slide.shapes, start=1):
            if shape.left < 0 or shape.top < 0 or shape.left + shape.width > slide_width or shape.top + shape.height > slide_height:
                errors.append(f"slide-{slide_number}-shape-{shape_number}-outside-canvas")
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text or ""
            text_chars += len(text.strip())
            errors.extend(
                f"slide-{slide_number}-placeholder:{hit}" for hit in _placeholder_hits(text)
            )
    if not presentation.slides:
        errors.append("PPTX-without-slides")
    if text_chars == 0:
        errors.append("PPTX-without-extractable-text")
    return {"format": "pptx", "slides": len(presentation.slides), "text_chars": text_chars, "errors": errors}


def validate_docx(path: Path) -> dict[str, Any]:
    from docx import Document

    document = Document(str(path))
    text = "\n".join(paragraph.text for paragraph in document.paragraphs)
    errors = [f"placeholder:{hit}" for hit in _placeholder_hits(text)]
    if not text.strip():
        errors.append("DOCX-without-extractable-text")
    with zipfile.ZipFile(path) as archive:
        if "word/document.xml" not in archive.namelist():
            errors.append("DOCX-missing-document-xml")
    return {"format": "docx", "paragraphs": len(document.paragraphs), "text_chars": len(text.strip()), "errors": errors}


def validate_file(path: Path) -> dict[str, Any]:
    suffix = path.suffix.casefold()
    if suffix == ".pdf":
        result = validate_pdf(path)
    elif suffix == ".pptx":
        result = validate_pptx(path)
    elif suffix == ".docx":
        result = validate_docx(path)
    else:
        return {"file": str(path), "format": suffix.lstrip("."), "errors": ["unsupported-format"]}
    return {"file": str(path), **result}


def validate_path(path: Path) -> dict[str, Any]:
    files = [path] if path.is_file() else sorted(item for item in path.rglob("*") if item.is_file())
    results = [validate_file(item) for item in files if item.suffix.casefold() in {".pdf", ".pptx", ".docx"}]
    errors = [f"{item['file']}: {error}" for item in results for error in item.get("errors", [])]
    return {"status": "passed" if not errors and results else "failed", "files": results, "errors": errors}


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--path", type=Path, required=True)
    parser.add_argument("--report", type=Path)
    args = parser.parse_args()
    report = validate_path(args.path)
    if args.report:
        args.report.parent.mkdir(parents=True, exist_ok=True)
        args.report.write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(report, ensure_ascii=False, indent=2))
    return 0 if report["status"] == "passed" else 1


if __name__ == "__main__":
    raise SystemExit(main())
