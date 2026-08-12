from __future__ import annotations

from pathlib import Path

from docx import Document
from pptx import Presentation
from pptx.util import Inches

from scripts.validate_exports import validate_path


def test_export_validator_accepts_minimal_valid_office_containers(tmp_path: Path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "Syntetisk læringsressurs"
    presentation.save(tmp_path / "lesson.pptx")

    document = Document()
    document.add_heading("Syntetisk læringsressurs", level=1)
    document.save(tmp_path / "lesson.docx")

    report = validate_path(tmp_path)
    assert report["status"] == "passed"
    assert {item["format"] for item in report["files"]} == {"pptx", "docx"}


def test_export_validator_rejects_internal_placeholder(tmp_path: Path):
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    box.text_frame.text = "TODO: legg inn kilde"
    presentation.save(tmp_path / "broken.pptx")

    report = validate_path(tmp_path)
    assert report["status"] == "failed"
    assert any("placeholder" in error for error in report["errors"])
