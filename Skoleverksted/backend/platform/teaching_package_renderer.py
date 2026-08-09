"""Deterministic renderers for TeachingPackage artifacts.

The renderer intentionally has no model dependency. Generation and truth
verification are separate concerns, so a teacher can still inspect a partial
package when one output format fails.
"""

from __future__ import annotations

import hashlib
import io
import re
from pathlib import Path
from typing import Iterable

from .models import TeachingArtifact, TeachingPackage, TruthSource


MIME_TYPES = {
    "pdf": "application/pdf",
    "docx": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "pptx": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}


def safe_teaching_filename(value: str, suffix: str) -> str:
    stem = re.sub(r"[^a-zA-Z0-9æøåÆØÅ_-]+", "-", value).strip("-").lower()
    return f"{stem[:100] or 'undervisningspakke'}.{suffix}"


def digest_bytes(value: bytes) -> str:
    return hashlib.sha256(value).hexdigest()


def _escape_typst(value: str) -> str:
    value = str(value or "").replace("\\", "\\\\")
    return re.sub(r"([#\[\]$<>@*_~`])", r"\\\1", value)


def _clean(value: str) -> str:
    value = re.sub(r"<[^>]+>", "", str(value or ""))
    value = re.sub(r"\*\*(.*?)\*\*", r"\1", value)
    value = re.sub(r"\[(.*?)\]\((https?://[^)]+)\)", r"\1 (\2)", value)
    return re.sub(r"\s+", " ", value).strip()


def _markdown_blocks(markdown: str) -> list[tuple[str, str]]:
    blocks: list[tuple[str, str]] = []
    current: list[str] = []
    heading = ""

    def flush() -> None:
        nonlocal current
        if current:
            text = "\n".join(current).strip()
            if text:
                blocks.append((heading, text))
        current = []

    for raw in str(markdown or "").replace("\\n", "\n").splitlines():
        match = re.match(r"^#{1,4}\s+(.+)$", raw.strip())
        if match:
            flush()
            heading = _clean(match.group(1))
        elif raw.strip():
            current.append(raw.strip())
    flush()
    return blocks


def _source_lines(sources: Iterable[TruthSource]) -> list[str]:
    values = []
    for source in sources:
        label = source.title or source.url
        values.append(f"- {_escape_typst(label)}: {_escape_typst(source.url)}")
    return values or ["- Ingen konkrete lærer- eller groundingkilder er registrert."]


def build_teaching_typst(package: TeachingPackage, artifact: TeachingArtifact) -> str:
    blocks = _markdown_blocks(artifact.content_markdown)
    body: list[str] = []
    for heading, text in blocks:
        if heading:
            body.append(f"== {_escape_typst(heading)}")
        for line in text.splitlines():
            stripped = line.strip()
            if re.match(r"^[-*]\s+", stripped):
                body.append(f"- {_escape_typst(_clean(re.sub(r'^[-*]\\s+', '', stripped)))}")
            elif re.match(r"^\d+[.)]\s+", stripped):
                body.append(f"+ {_escape_typst(_clean(re.sub(r'^\\d+[.)]\\s+', '', stripped)))}")
            else:
                body.append(_escape_typst(_clean(stripped)))
            body.append("#v(4pt)")
    goals = "\n".join(f"- {_escape_typst(item)}" for item in package.plan.learning_goals) or "- Målene må konkretiseres før bruk."
    concepts = ", ".join(_clean(item) for item in package.plan.key_concepts) or "Ikke registrert"
    sources = "\n".join(_source_lines(artifact.sources or package.plan.sources))
    return f'''#set page(paper: "a4", margin: (x: 20mm, y: 18mm))
#set text(font: "Aptos", size: 10.5pt, lang: "nb")
#set par(leading: 0.65em, spacing: 0.7em)
#show heading.where(level: 1): set text(size: 19pt, weight: "bold", fill: rgb("#183B56"))
#show heading.where(level: 2): set text(size: 14pt, weight: "bold", fill: rgb("#235F72"))

#text(size: 9pt, weight: "bold", fill: rgb("#235F72"))[SKOLEVERKSTED · {_escape_typst(package.subject.upper())}]
#v(8pt)
#text(size: 25pt, weight: "bold", fill: rgb("#102A43"))[{_escape_typst(artifact.title)}]
#v(3pt)
#text(size: 11pt, fill: rgb("#52606D"))[{_escape_typst(package.level)} · {_escape_typst(package.plan.audience)} · {_escape_typst(package.plan.theme)}]
#v(14pt)
#block(fill: rgb("#E8F1F5"), inset: 10pt, radius: 4pt)[
  #strong[Læringsmål]
  #v(3pt)
  {goals}
  #v(4pt)
  #strong[Sentrale begreper:] {_escape_typst(concepts)}
]
#v(12pt)
{chr(10).join(body)}

#v(12pt)
== Kilder og kvalitetsgrunnlag
Faktapass og kildestatus gjelder den eksakte innholdsrevisjonen som er lagret i pakken. Læreren må lese kildene og vurdere faglig egnethet før utdeling.
#v(5pt)
{sources}
'''


def render_pdf(package: TeachingPackage, artifact: TeachingArtifact) -> bytes:
    from VGS_KI.backend.pdf_service import compile_typst

    return compile_typst(build_teaching_typst(package, artifact))


def _add_docx_markdown(document, markdown: str) -> None:
    for heading, text in _markdown_blocks(markdown):
        if heading:
            document.add_heading(heading, level=2)
        for line in text.splitlines():
            stripped = line.strip()
            if re.match(r"^[-*]\s+", stripped):
                document.add_paragraph(re.sub(r"^[-*]\s+", "", stripped), style="List Bullet")
            elif re.match(r"^\d+[.)]\s+", stripped):
                document.add_paragraph(re.sub(r"^\d+[.)]\s+", "", stripped), style="List Number")
            else:
                document.add_paragraph(_clean(stripped))


def render_docx(package: TeachingPackage, artifact: TeachingArtifact) -> bytes:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt, RGBColor

    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.75)
    section.bottom_margin = Inches(0.75)
    section.left_margin = Inches(0.85)
    section.right_margin = Inches(0.85)
    normal = document.styles["Normal"]
    normal.font.name = "Aptos"
    normal.font.size = Pt(10.5)
    title = document.add_heading(artifact.title, level=0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for run in title.runs:
        run.font.color.rgb = RGBColor(0x10, 0x2A, 0x43)
    subtitle = document.add_paragraph(f"{package.subject} · {package.level} · {package.plan.audience}")
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    subtitle.runs[0].font.color.rgb = RGBColor(0x52, 0x60, 0x6D)
    document.add_paragraph(package.plan.overview or f"Tema: {package.plan.theme}")
    document.add_heading("Læringsmål", level=1)
    for goal in package.plan.learning_goals or ["Målene må konkretiseres før bruk."]:
        document.add_paragraph(goal, style="List Bullet")
    document.add_heading("Innhold", level=1)
    _add_docx_markdown(document, artifact.content_markdown)
    document.add_heading("Kilder og kvalitetsgrunnlag", level=1)
    document.add_paragraph(
        "Faktapass og kildestatus gjelder den eksakte innholdsrevisjonen som er lagret i pakken. "
        "Læreren må lese kildene og vurdere faglig egnethet før utdeling."
    )
    for source in artifact.sources or package.plan.sources:
        paragraph = document.add_paragraph(style="List Bullet")
        paragraph.add_run(f"{source.title}: {source.url}")
    footer = section.footer.paragraphs[0]
    footer.alignment = WD_ALIGN_PARAGRAPH.CENTER
    footer_run = footer.add_run("Skoleverksted · TeachingPackage")
    footer_run.font.size = Pt(8)
    footer_run.font.color.rgb = RGBColor(0x82, 0x8B, 0x98)
    stream = io.BytesIO()
    document.save(stream)
    return stream.getvalue()


def _presentation_sections(markdown: str) -> list[tuple[str, list[str]]]:
    sections: list[tuple[str, list[str]]] = []
    current_title = ""
    current_lines: list[str] = []
    for raw in str(markdown or "").replace("\\n", "\n").splitlines():
        match = re.match(r"^#{1,4}\s+(.+)$", raw.strip())
        if match:
            if current_title:
                sections.append((current_title, current_lines))
            current_title = _clean(match.group(1))
            current_lines = []
        elif raw.strip():
            current_lines.append(_clean(raw.strip()))
    if current_title:
        sections.append((current_title, current_lines))
    return sections


def render_pptx(package: TeachingPackage, artifact: TeachingArtifact) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
    from pptx.util import Inches, Pt

    presentation = Presentation()
    presentation.slide_width = Inches(13.333)
    presentation.slide_height = Inches(7.5)
    sections = _presentation_sections(artifact.content_markdown)
    if not sections:
        sections = [(artifact.title, [package.plan.overview or package.plan.theme])]
    sources = artifact.sources or package.plan.sources
    for index, (title, lines) in enumerate(sections):
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = RGBColor(0xF7, 0xFA, 0xFC)
        band = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(0.16))
        band.fill.solid()
        band.fill.fore_color.rgb = RGBColor(0x23, 0x5F, 0x72)
        band.line.fill.background()
        title_box = slide.shapes.add_textbox(Inches(0.7), Inches(0.55), Inches(11.9), Inches(0.85))
        title_frame = title_box.text_frame
        title_frame.clear()
        title_frame.word_wrap = False
        title_para = title_frame.paragraphs[0]
        title_para.text = title[:92]
        title_para.font.size = Pt(34 if index else 38)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(0x10, 0x2A, 0x43)
        title_para.alignment = PP_ALIGN.LEFT
        body_box = slide.shapes.add_textbox(Inches(0.85), Inches(1.65), Inches(11.55), Inches(4.9))
        frame = body_box.text_frame
        frame.clear()
        frame.word_wrap = True
        frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        display_lines = lines[:8] or [package.plan.overview or package.plan.theme]
        for line_index, line in enumerate(display_lines):
            paragraph = frame.paragraphs[0] if line_index == 0 else frame.add_paragraph()
            paragraph.text = line[:260]
            paragraph.font.size = Pt(22)
            paragraph.font.color.rgb = RGBColor(0x24, 0x3B, 0x53)
            paragraph.space_after = Pt(12)
            if line.startswith(("- ", "* ")):
                paragraph.text = f"• {line[2:]}"
                paragraph.level = 0
                paragraph.font.size = Pt(21)
        footer = slide.shapes.add_textbox(Inches(0.85), Inches(6.85), Inches(11.5), Inches(0.28))
        footer.text_frame.text = f"{package.plan.theme} · {index + 1}/{len(sections)}"
        footer.text_frame.paragraphs[0].font.size = Pt(11)
        footer.text_frame.paragraphs[0].font.color.rgb = RGBColor(0x6B, 0x7C, 0x8F)
        notes = slide.notes_slide.notes_text_frame
        notes.text = (
            "Forslag til gjennomføring: Bruk lysbildet som utgangspunkt for spørsmål og elevaktivitet. "
            "Tilpass tempoet til klassen.\n\n[Sources]\n"
            + "\n".join(f"- {source.title}: {source.url}" for source in sources)
            if sources
            else "Forslag til gjennomføring: Bruk lysbildet som utgangspunkt for spørsmål og elevaktivitet.\n\n[Sources]\n- Ingen konkrete kilder registrert."
        )
    stream = io.BytesIO()
    presentation.save(stream)
    return stream.getvalue()


def render_artifact(package: TeachingPackage, artifact: TeachingArtifact) -> dict[str, bytes]:
    if artifact.artifact_type == "presentation":
        return {"pptx": render_pptx(package, artifact)}
    return {"pdf": render_pdf(package, artifact), "docx": render_docx(package, artifact)}


def artifact_file_metadata(
    package: TeachingPackage,
    artifact: TeachingArtifact,
    rendered: dict[str, bytes],
) -> list[dict[str, object]]:
    return [
        {
            "format": fmt,
            "filename": safe_teaching_filename(artifact.title, fmt),
            "mime_type": MIME_TYPES[fmt],
            "size_bytes": len(content),
            "digest": digest_bytes(content),
            "storage_key": f"{package.id}/r{package.package_revision}/{artifact.id}.{fmt}",
            "package_revision": package.package_revision,
        }
        for fmt, content in rendered.items()
    ]


def rendered_paths(base_dir: Path, package: TeachingPackage, artifact: TeachingArtifact) -> list[Path]:
    return [base_dir / file.storage_key for file in artifact.files]
