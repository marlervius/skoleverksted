"""
PowerPoint Export for MateMaTeX.
Export content to PowerPoint presentations.
"""

import re
from pathlib import Path
from typing import Optional
from dataclasses import dataclass


def is_pptx_available() -> bool:
    """Check if python-pptx is installed."""
    try:
        import pptx
        return True
    except ImportError:
        return False


@dataclass
class SlideContent:
    """Content for a single slide."""
    title: str
    content: list[str]
    slide_type: str  # "title", "content", "exercise", "solution"
    notes: Optional[str] = None


def parse_latex_to_slides(latex_content: str) -> list[SlideContent]:
    """
    Parse LaTeX content into slide content.
    
    Args:
        latex_content: The LaTeX source code.
    
    Returns:
        List of SlideContent objects.
    """
    slides = []
    
    # Extract document title
    title_match = re.search(r'\\title\{([^}]+)\}', latex_content)
    doc_title = title_match.group(1) if title_match else "Matematikk"
    
    # Add title slide
    slides.append(SlideContent(
        title=doc_title,
        content=["Generert med MateMaTeX"],
        slide_type="title"
    ))
    
    # Find sections
    section_pattern = r'\\section\*?\{([^}]+)\}'
    sections = re.split(section_pattern, latex_content)
    
    current_section = None
    
    for i, part in enumerate(sections):
        if i % 2 == 1:  # This is a section title
            current_section = part.strip()
            continue
        
        if not current_section:
            continue
        
        # Extract content from this section
        content_items = []
        
        # Find definitions
        def_pattern = r'\\begin\{definition\}(.*?)\\end\{definition\}'
        for match in re.finditer(def_pattern, part, re.DOTALL):
            text = clean_latex_for_pptx(match.group(1))
            if text:
                content_items.append(f"📘 {text}")
        
        # Find examples
        example_pattern = r'\\begin\{example\}(.*?)\\end\{example\}'
        for match in re.finditer(example_pattern, part, re.DOTALL):
            text = clean_latex_for_pptx(match.group(1))
            if text:
                content_items.append(f"💡 {text}")
        
        # Find tasks/exercises
        task_pattern = r'\\begin\{taskbox\}\{([^}]*)\}(.*?)\\end\{taskbox\}'
        exercises = []
        for match in re.finditer(task_pattern, part, re.DOTALL):
            task_title = match.group(1).strip()
            task_content = clean_latex_for_pptx(match.group(2))
            if task_content:
                exercises.append(f"{task_title}: {task_content[:100]}...")
        
        # Create section slide
        if content_items:
            slides.append(SlideContent(
                title=current_section,
                content=content_items[:5],  # Max 5 items per slide
                slide_type="content"
            ))
        
        # Create exercise slides (group 3-4 per slide)
        for j in range(0, len(exercises), 3):
            chunk = exercises[j:j+3]
            slides.append(SlideContent(
                title=f"Oppgaver - {current_section}",
                content=chunk,
                slide_type="exercise"
            ))
    
    return slides


def clean_latex_for_pptx(text: str) -> str:
    """Clean LaTeX markup for PowerPoint display."""
    # Remove common LaTeX commands
    text = re.sub(r'\\textbf\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\textit\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\emph\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\underline\{([^}]*)\}', r'\1', text)
    
    # Convert math to text representation
    text = re.sub(r'\$([^$]+)\$', r'[\1]', text)
    text = re.sub(r'\\\[([^\]]+)\\\]', r'[\1]', text)
    
    # Remove other commands
    text = re.sub(r'\\[a-zA-Z]+\{[^}]*\}', '', text)
    text = re.sub(r'\\[a-zA-Z]+', '', text)
    
    # Clean up whitespace
    text = re.sub(r'\s+', ' ', text).strip()
    
    return text


def create_pptx(
    slides: list[SlideContent],
    output_path: str,
    theme_color: str = "#f0b429"
) -> Optional[str]:
    """
    Create a PowerPoint presentation.
    
    Args:
        slides: List of slide content.
        output_path: Path to save the PPTX file.
        theme_color: Primary theme color (hex).
    
    Returns:
        Path to the created file, or None if failed.
    """
    if not is_pptx_available():
        return None
    
    try:
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RgbColor
        from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
        
        # Create presentation
        prs = Presentation()
        prs.slide_width = Inches(13.333)  # 16:9 widescreen
        prs.slide_height = Inches(7.5)
        
        # Get layouts
        title_layout = prs.slide_layouts[6]  # Blank
        content_layout = prs.slide_layouts[6]  # Blank
        
        for slide_content in slides:
            slide = prs.slides.add_slide(content_layout)
            
            if slide_content.slide_type == "title":
                # Title slide
                # Add title
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(2.5), Inches(12.333), Inches(1.5)
                )
                tf = title_box.text_frame
                tf.word_wrap = True
                p = tf.paragraphs[0]
                p.text = slide_content.title
                p.font.size = Pt(44)
                p.font.bold = True
                p.alignment = PP_ALIGN.CENTER
                
                # Subtitle
                if slide_content.content:
                    sub_box = slide.shapes.add_textbox(
                        Inches(0.5), Inches(4.2), Inches(12.333), Inches(0.75)
                    )
                    tf = sub_box.text_frame
                    p = tf.paragraphs[0]
                    p.text = slide_content.content[0]
                    p.font.size = Pt(24)
                    p.alignment = PP_ALIGN.CENTER
            
            else:
                # Content slide
                # Add title
                title_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8)
                )
                tf = title_box.text_frame
                p = tf.paragraphs[0]
                p.text = slide_content.title
                p.font.size = Pt(32)
                p.font.bold = True
                
                # Add content
                content_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(1.5), Inches(12.333), Inches(5.5)
                )
                tf = content_box.text_frame
                tf.word_wrap = True
                
                for i, item in enumerate(slide_content.content):
                    if i == 0:
                        p = tf.paragraphs[0]
                    else:
                        p = tf.add_paragraph()
                    
                    p.text = f"• {item}"
                    p.font.size = Pt(20)
                    p.space_before = Pt(12)
                    p.space_after = Pt(6)
            
            # Add notes if present
            if slide_content.notes:
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = slide_content.notes
        
        # Save presentation
        output_path = Path(output_path)
        output_path.parent.mkdir(parents=True, exist_ok=True)
        prs.save(str(output_path))
        
        return str(output_path)
    
    except Exception as e:
        print(f"Error creating PPTX: {e}")
        return None


def latex_to_pptx(latex_content: str, output_path: str) -> Optional[str]:
    """
    Convert LaTeX content to PowerPoint presentation.
    
    Args:
        latex_content: The LaTeX source code.
        output_path: Path to save the PPTX file.
    
    Returns:
        Path to the created file, or None if failed.
    """
    slides = parse_latex_to_slides(latex_content)
    return create_pptx(slides, output_path)


def get_pptx_preview(slides: list[SlideContent]) -> str:
    """Generate HTML preview of the presentation structure."""
    html = '<div style="display: flex; flex-wrap: wrap; gap: 0.5rem;">'
    
    for i, slide in enumerate(slides):
        bg_color = {
            "title": "rgba(240,180,41,0.2)",
            "content": "rgba(59,130,246,0.2)",
            "exercise": "rgba(16,185,129,0.2)",
            "solution": "rgba(139,92,246,0.2)",
        }.get(slide.slide_type, "rgba(255,255,255,0.1)")
        
        icon = {
            "title": "🎬",
            "content": "📖",
            "exercise": "✍️",
            "solution": "🔑",
        }.get(slide.slide_type, "📄")
        
        html += f"""
        <div style="
            width: 120px;
            height: 70px;
            background: {bg_color};
            border: 1px solid rgba(255,255,255,0.2);
            border-radius: 4px;
            padding: 0.5rem;
            font-size: 0.6rem;
            overflow: hidden;
        ">
            <div style="display: flex; align-items: center; gap: 0.25rem; margin-bottom: 0.25rem;">
                <span>{icon}</span>
                <span style="color: #e2e8f0; font-weight: 500;">Lysbilde {i+1}</span>
            </div>
            <div style="color: #9090a0; font-size: 0.55rem; line-height: 1.2;">
                {slide.title[:20]}{'...' if len(slide.title) > 20 else ''}
            </div>
        </div>
        """
    
    html += '</div>'
    return html
