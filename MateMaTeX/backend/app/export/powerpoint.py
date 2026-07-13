"""
LaTeX → PowerPoint (PPTX) conversion.

Each exercise becomes its own slide. Solutions go to speaker notes
or hidden slides.
"""

from __future__ import annotations

import io
import re

import structlog

logger = structlog.get_logger()


def _extract_exercises_for_slides(latex: str) -> list[dict]:
    """Extract exercises from LaTeX for slide-by-slide conversion."""
    exercises = []

    # Match taskbox environments
    pattern = re.compile(
        r'\\begin\{taskbox\}\{([^}]*)\}(.*?)\\end\{taskbox\}',
        re.DOTALL,
    )

    for match in pattern.finditer(latex):
        title = match.group(1).strip()
        body = match.group(2).strip()
        exercises.append({"title": title, "body": body})

    if not exercises:
        # Fallback: split by \section or double newlines
        sections = re.split(r'\\section\*?\{([^}]*)\}', latex)
        for i in range(1, len(sections), 2):
            title = sections[i].strip() if i < len(sections) else f"Slide {i}"
            body = sections[i + 1].strip() if i + 1 < len(sections) else ""
            if body:
                exercises.append({"title": title, "body": body})

    return exercises


def _simplify_latex_for_slide(text: str) -> str:
    """Strip LaTeX to readable text for PowerPoint slides."""
    # Remove environments
    text = re.sub(r'\\begin\{[^}]*\}(?:\{[^}]*\})?', '', text)
    text = re.sub(r'\\end\{[^}]*\}', '', text)

    # Simplify math
    text = re.sub(r'\$([^$]*)\$', r'\1', text)
    text = re.sub(r'\\frac\{([^}]*)\}\{([^}]*)\}', r'(\1)/(\2)', text)
    text = re.sub(r'\\sqrt\{([^}]*)\}', r'√(\1)', text)
    text = re.sub(r'\\cdot', r'·', text)
    text = re.sub(r'\\times', r'×', text)

    # Clean commands
    text = re.sub(r'\\textbf\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\textit\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\item', r'• ', text)
    text = re.sub(r'\\[a-zA-Z]+\*?\{([^}]*)\}', r'\1', text)
    text = re.sub(r'\\[a-zA-Z]+\*?', '', text)
    text = re.sub(r'[{}]', '', text)
    text = re.sub(r'\n{3,}', '\n\n', text)

    return text.strip()


def latex_to_pptx(
    latex_content: str,
    title: str = "Matematikk",
    solutions_as: str = "speaker_notes",
) -> bytes:
    """
    Convert LaTeX content to a PowerPoint presentation.

    Returns the PPTX file as bytes.

    Args:
        latex_content: Full LaTeX content
        title: Presentation title
        solutions_as: 'speaker_notes' or 'hidden_slides'
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.text import PP_ALIGN
    from pptx.dml.color import RGBColor

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Title slide
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Generert av MateMaTeX AI"

    # Extract exercises
    exercises = _extract_exercises_for_slides(latex_content)

    # Extract solutions section
    solutions: dict[int, str] = {}
    sol_match = re.search(
        r'\\section\*?\{Løsningsforslag\}(.*)',
        latex_content,
        re.DOTALL,
    )
    if sol_match:
        sol_text = sol_match.group(1)
        for m in re.finditer(
            r'\\textbf\{Oppgave\s*(\d+)\}(.*?)(?=\\textbf\{Oppgave|\Z)',
            sol_text,
            re.DOTALL,
        ):
            solutions[int(m.group(1))] = _simplify_latex_for_slide(m.group(2))

    for i, ex in enumerate(exercises, 1):
        # Exercise slide
        slide_layout = prs.slide_layouts[1]  # Title + content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = ex["title"]

        body = _simplify_latex_for_slide(ex["body"])
        tf = slide.placeholders[1].text_frame
        tf.text = body
        for paragraph in tf.paragraphs:
            paragraph.font.size = Pt(18)

        # Solution
        solution = solutions.get(i, "")
        if solution:
            if solutions_as == "speaker_notes":
                notes_slide = slide.notes_slide
                notes_slide.notes_text_frame.text = f"Løsning:\n{solution}"
            elif solutions_as == "hidden_slides":
                sol_slide_layout = prs.slide_layouts[1]
                sol_slide = prs.slides.add_slide(sol_slide_layout)
                sol_slide.shapes.title.text = f"Løsning — {ex['title']}"
                tf = sol_slide.placeholders[1].text_frame
                tf.text = solution
                for paragraph in tf.paragraphs:
                    paragraph.font.size = Pt(16)

    # Save to bytes
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
